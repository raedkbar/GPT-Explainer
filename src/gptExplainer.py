import json
import os
import asyncio
import openai
import pptx


API_KEY = "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"
PROMPT_INIT = "Explain the content of the following slide. " \
              "Write a response as if you were writing an article, and don't break the fourth wall! " \
              "Meaning, don't mention the words slide or presentation:"
MODEL_VERSION = "gpt-3.5-turbo"
TIMEOUT_SECONDS = 30  # Timeout value in seconds
MAX_RETRIES = 3  # Maximum number of retries for slide processing
UPLOADS_DIR = "../uploads"
OUTPUTS_DIR = "../outputs"


class SlideProcessingError(Exception):
    pass


class PresentationProcessingError(Exception):
    pass


class OpenAIError(Exception):
    pass


async def process_slide(slide_num, slide_text, retry_count=1):
    """
    Processes a slide by generating an explanation using OpenAI's ChatCompletion model.

    Args:
        slide_num (int): The slide number.
        slide_text (str): The text content of the slide.
        retry_count (int): The number of times the slide processing has been retried (default: 1).

    Returns:
        str: The generated explanation for the slide.

    Raises:
        SlideProcessingError: If the slide processing fails after the maximum number of retries.
    """
    try:
        print(f"Processing slide {slide_num}...\n")
        prompt = f"{PROMPT_INIT}\n{slide_text}\n"
        response = await openai.ChatCompletion.acreate(
            model=MODEL_VERSION,
            messages=[
                {"role": "system", "content": "You are a helpful assistant that explains PowerPoint slides:"},
                {"role": "user", "content": prompt}
            ]
        )
        explanation = response.choices[0].message.content
        return explanation
    except Exception as e:
        if retry_count < MAX_RETRIES:
            print(f"Error occurred while processing slide {slide_num}: {str(e)}\n")
            print(f"Retrying to process slide {slide_num}. Retry count: {retry_count + 1}\n")
            return await process_slide(slide_num, slide_text, retry_count=retry_count + 1)
        else:
            raise SlideProcessingError(f"Failed to process slide {slide_num} after {MAX_RETRIES} retries: {str(e)}") from e


def extract_slide_text(slide):
    """
    Extracts the text content from a slide.

    Args:
        slide: The slide object.

    Returns:
        str: The extracted text content of the slide.
    """
    slide_text = ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    slide_text += run.text
    return slide_text.strip()


async def process_presentation(presentation_path):
    """
    Processes a PowerPoint presentation by extracting slide texts and generating explanations for each slide.

    Args:
        presentation_path (str): The path to the PowerPoint presentation file.

    Returns:
        list: A list of dictionaries containing the slide number and its corresponding explanation.

    Raises:
        PresentationProcessingError: If the presentation processing fails.
    """
    try:
        presentation = pptx.Presentation(presentation_path)
        slides = [extract_slide_text(slide) for slide in presentation.slides]

        explanations = []

        tasks = []
        for i, slide_text in enumerate(slides, start=1):
            if slide_text:
                task = asyncio.create_task(process_slide(i, slide_text))
                tasks.append(task)

        results = await asyncio.gather(*tasks)

        for i, explanation in enumerate(results, start=1):
            explanations.append({"slide": i, "explanation": explanation})

        return explanations
    except Exception as e:
        raise PresentationProcessingError(f"Failed to process presentation: {str(e)}") from e


async def process_file(file_path):
    """
    Processes a file by generating explanations for each slide in the PowerPoint presentation.

    Args:
        file_path (str): The path to the PowerPoint presentation file.
    """
    try:
        output_file = os.path.join(OUTPUTS_DIR, os.path.splitext(os.path.basename(file_path))[0] + ".json")

        openai.api_key = API_KEY

        print(f"Processing file: {file_path}\n")

        explanations = await process_presentation(file_path)

        # Save explanations to a JSON file
        with open(output_file, "w") as f:
            json.dump(explanations, f, indent=4)

        print(f"Explanations saved successfully for file: {file_path}")
    except (OpenAIError, PresentationProcessingError, SlideProcessingError) as e:
        print(f"Error occurred while processing file {file_path}: {str(e)}")


async def process_uploads_folder():
    """
    Processes the files in the 'uploads' folder and generates explanations for each PowerPoint presentation.

    The function runs indefinitely, scanning the 'uploads' folder every few seconds and processing new files.
    """
    while True:
        print("Scanning uploads folder...\n")
        files = os.listdir(UPLOADS_DIR)

        for file in files:
            file_path = os.path.join(UPLOADS_DIR, file)
            if file.endswith(".pptx") and not file.startswith("processed_"):
                await process_file(file_path)
                os.rename(file_path, os.path.join(UPLOADS_DIR, "processed_" + file))

        print("Sleeping for 10 seconds...\n")
        await asyncio.sleep(10)


if __name__ == "__main__":
    asyncio.run(process_uploads_folder())
