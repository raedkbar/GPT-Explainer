import asyncio
import pptx
from gpt_processing import process_slide
from exceptions import PresentationProcessingError


async def extract_slide_text(slide):
    """
    Extracts the text content from a slide.

    Args:
        slide: The slide object.

    Returns:
        str: The extracted text content of the slide.
    """
    slide_text = ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    slide_text += run.text
    return slide_text.strip()


async def process_presentation(presentation_path):
    """
    Processes a PowerPoint presentation by extracting slide texts and generating explanations for each slide.

    Args:
        presentation_path (str): The path to the PowerPoint presentation file.

    Returns:
        list: A list of dictionaries containing the slide number and its corresponding explanation.

    Raises:
        PresentationProcessingError: If the presentation processing fails.
    """
    try:
        presentation = pptx.Presentation(presentation_path)
        slides = [await extract_slide_text(slide) for slide in presentation.slides]

        explanations = []

        tasks = []
        for i, slide_text in enumerate(slides, start=1):
            if slide_text:
                task = asyncio.create_task(process_slide(i, slide_text))
                tasks.append(task)

        results = await asyncio.gather(*tasks)

        for i, explanation in enumerate(results, start=1):
            explanations.append({"slide": i, "explanation": explanation})

        return explanations
    except Exception as e:
        raise PresentationProcessingError(str(e)) from e
