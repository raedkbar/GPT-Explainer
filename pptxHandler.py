import asyncio
import os
import json
import pptx
import openai
import argparse

API_KEY = "API_KEY"
PROMPT_INIT = "Explain the content of the following slide. " \
              "Write a response as if you were writing an article, and don't break the fourth wall! " \
              "Meaning, don't mention the words slide or presentation:"
TIMEOUT_SECONDS = 30  # Timeout value in seconds
MAX_RETRIES = 3  # Maximum number of retries for slide processing


class SlideProcessingError(Exception):
    pass


class PresentationProcessingError(Exception):
    pass


class OpenAIError(Exception):
    pass


async def process_slide(slide_num, slide_text, retry_count=1):
    """
    Processes a slide by generating an explanation using OpenAI's ChatCompletion model.

    Args:
        slide_num (int): The slide number.
        slide_text (str): The text content of the slide.
        retry_count (int): The number of times the slide processing has been retried (default: 1).

    Returns:
        str: The generated explanation for the slide.

    Raises:
        SlideProcessingError: If the slide processing fails after the maximum number of retries.
    """
    try:
        print(f"Processing slide {slide_num}...\n")
        prompt = f"{PROMPT_INIT}\n{slide_text}\n"
        response = await openai.ChatCompletion.acreate(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are a helpful assistant that explains PowerPoint slides:"},
                {"role": "user", "content": prompt}
            ]
        )
        explanation = response.choices[0].message.content
        return explanation
    except Exception as e:
        if retry_count < MAX_RETRIES:
            print(f"Error occurred while processing slide {slide_num}: {str(e)}\n")
            print(f"Retrying to process slide {slide_num}. Retry count: {retry_count + 1}\n")
            return await process_slide(slide_num, slide_text, retry_count=retry_count + 1)
        else:
            raise SlideProcessingError(f"Failed to process slide after {MAX_RETRIES} retries: {str(e)}") from e


def extract_slide_text(slide):
    """
    Extracts the text content from a slide.

    Args:
        slide: The slide object.

    Returns:
        str: The extracted text content of the slide.
    """
    slide_text = ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    slide_text += run.text
    return slide_text.strip()


async def process_presentation(presentation_path):
    """
    Processes a PowerPoint presentation by extracting slide texts and generating explanations for each slide.

    Args:
        presentation_path (str): The path to the PowerPoint presentation file.

    Returns:
        list: A list of dictionaries containing the slide number and its corresponding explanation.

    Raises:
        PresentationProcessingError: If the presentation processing fails.
    """
    try:
        presentation = pptx.Presentation(presentation_path)
        slides = [extract_slide_text(slide) for slide in presentation.slides]

        explanations = []

        for i, slide_text in enumerate(slides, start=1):
            if slide_text:
                explanation = await process_slide(i, slide_text)
            else:
                explanation = None

            explanations.append({"slide": i, "explanation": explanation})

        return explanations
    except Exception as e:
        raise PresentationProcessingError(f"Failed to process presentation: {str(e)}") from e


async def main(presentation_path):
    """
    Main function to process the PowerPoint presentation and save the explanations to a JSON file.

    Args:
        presentation_path (str): The path to the PowerPoint presentation file.
    """
    try:
        output_file = os.path.splitext(presentation_path)[0] + ".json"

        openai.api_key = API_KEY

        print("Processing file...\n")

        explanations = await process_presentation(presentation_path)

        # Save explanations to a JSON file
        with open(output_file, "w") as f:
            json.dump(explanations, f, indent=4)

        print("Explanations saved successfully.")
    except (OpenAIError, PresentationProcessingError, SlideProcessingError) as e:
        print(f"Error occurred: {str(e)}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Process PowerPoint presentation slides.")
    parser.add_argument("presentation_path", type=str, help="Path to the PowerPoint presentation file.")
    args = parser.parse_args()

    asyncio.run(main(args.presentation_path))
