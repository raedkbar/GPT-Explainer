import asyncio
import json
import os
import pptx
from datetime import datetime

from pkg.db_util.db_session import create_session
from pkg.explainer.exceptions import PresentationProcessingError
from pkg.explainer.gpt_processing import process_slide
from pkg.db_util.ORM import Upload

UPLOADS_DIR = "../uploads"
OUTPUTS_DIR = "../outputs"

session = create_session()


async def extract_slide_text(slide):
    """
    Extracts the text content from a slide.

    Args:
        slide: The slide object.

    Returns:
        str: The extracted text content of the slide.
    """
    slide_text = ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    slide_text += run.text
    return slide_text.strip()


async def process_presentation(upload_id):
    """
    Processes a PowerPoint presentation by extracting slide texts and generating explanations for each slide.

    Args:
        upload_id (int): The ID of the upload record.

    Raises:
        PresentationProcessingError: If the presentation processing fails.
    """
    try:
        upload = session.get(Upload, upload_id)

        if upload is None:
            print(f"\nUpload with ID '{upload_id}' not found")
            return

        presentation_path = f"{UPLOADS_DIR}/{upload.uid}.pptx"
        presentation = pptx.Presentation(presentation_path)
        slides = [await extract_slide_text(slide) for slide in presentation.slides]

        explanations = []

        tasks = []
        for i, slide_text in enumerate(slides, start=1):
            if slide_text:
                task = asyncio.create_task(process_slide(i, slide_text))
                tasks.append(task)

        results = await asyncio.gather(*tasks)

        for i, explanation in enumerate(results, start=1):
            explanations.append({"slide": i, "explanation": explanation})

        output_file = os.path.join(OUTPUTS_DIR, os.path.splitext(os.path.basename(upload.uid))[0] + ".json")

        # Save explanations to a JSON file
        with open(output_file, "w") as f:
            json.dump(explanations, f, indent=4)

        # Update the upload record with explanations and finish time
        upload.status = "completed"
        upload.finish_time = datetime.now()
        session.commit()

        print(f"Explanations saved successfully for upload ID: {upload_id}\n")
    except Exception as e:
        raise PresentationProcessingError(str(e)) from e
